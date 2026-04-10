from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "deliverables" / "dashboard-proposal-2026-04-15.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(17, 37, 66)
BLUE = RGBColor(43, 108, 176)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(180, 83, 9)
RED = RGBColor(185, 28, 28)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(107, 114, 128)
LIGHT = RGBColor(243, 244, 246)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(209, 213, 219)
TITLE_FONT = "Arial"
BODY_FONT = "Arial"


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    color = WHITE if dark else NAVY
    sub_color = RGBColor(226, 232, 240) if dark else MUTED

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(11.8), Inches(0.75))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = color
    r.font.name = TITLE_FONT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.22), Inches(11.5), Inches(0.45))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(11)
        r.font.color.rgb = sub_color
        r.font.name = BODY_FONT


def add_footer(slide, left: str, right: str, dark: bool = False) -> None:
    footer_color = RGBColor(203, 213, 225) if dark else MUTED
    line_color = RGBColor(148, 163, 184) if dark else BORDER
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(7.0), Inches(11.8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = line_color
    line.line.fill.background()

    left_box = slide.shapes.add_textbox(Inches(0.78), Inches(7.05), Inches(4.5), Inches(0.22))
    p = left_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = left
    r.font.size = Pt(9)
    r.font.color.rgb = footer_color
    r.font.name = BODY_FONT

    right_box = slide.shapes.add_textbox(Inches(9.6), Inches(7.05), Inches(2.9), Inches(0.22))
    p = right_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = right
    r.font.size = Pt(9)
    r.font.color.rgb = footer_color
    r.font.name = BODY_FONT


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 18, color: RGBColor = TEXT) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(10)
        p.bullet = True
        r = p.add_run()
        r.text = item
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.name = BODY_FONT


def add_callout(slide, title: str, body: str, left: float, top: float, width: float, height: float, accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = BORDER

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(0.12), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.18), Inches(width - 0.4), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = TITLE_FONT

    body_box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.58), Inches(width - 0.4), Inches(height - 0.75))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = body
    r.font.size = Pt(11)
    r.font.color.rgb = TEXT
    r.font.name = BODY_FONT


def add_panel(slide, label: str, body: str, left: float, top: float, width: float, height: float, fill_color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color

    label_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.18), Inches(width - 0.35), Inches(0.3))
    p = label_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = label
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = TITLE_FONT

    body_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.55), Inches(width - 0.35), Inches(height - 0.7))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(10.5)
        r.font.color.rgb = WHITE
        r.font.name = BODY_FONT


def add_banner(slide, text: str, left: float, top: float, width: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    box = slide.shapes.add_textbox(Inches(left + 0.08), Inches(top + 0.07), Inches(width - 0.16), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = TITLE_FONT


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_background(slide, NAVY)
    add_title(
        slide,
        "Ride YourWay Dashboard MVP Proposal",
        "Executive presentation | prepared April 15, 2026 | proposed delivery April 16, 2026",
        dark=True,
    )
    add_banner(slide, "DECISION DASHBOARD MVP, NOT A PRODUCTION PLATFORM", 0.78, 1.75, 6.0, TEAL)
    add_callout(
        slide,
        "What we are committing to",
        "Deliver a market-readiness dashboard MVP and explanation deck by Thursday, April 16, 2026, using only the frozen workspace data, rules, and assumptions.",
        0.8,
        2.35,
        5.6,
        1.55,
        TEAL,
    )
    add_callout(
        slide,
        "Why the commitment is credible",
        "The business question is frozen, the analytical base is already extracted, and the unresolved formulas are known well enough to ship with confidence labels rather than fake certainty.",
        6.6,
        2.35,
        5.8,
        1.55,
        BLUE,
    )
    add_bullets(
        slide,
        [
            "Objective: determine if a candidate market can meet nine launch gates and sustain a 25% operating margin.",
            "Approach: rules-based engine with assumptions disclosed, not a machine-learning forecast.",
            "Output: readiness scorecard, projected margin, risk flags, mitigation levers, and a short supporting deck.",
        ],
        0.95,
        4.35,
        11.0,
        1.8,
        font_size=18,
        color=WHITE,
    )
    add_footer(slide, "Ride YourWay Dashboard MVP Proposal", "1", dark=True)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    add_title(slide, "What Will Be Delivered Tomorrow", "A tight MVP scope keeps the promise achievable and defensible")
    add_panel(slide, "1. Dashboard MVP", "Accepts the intake example\nReturns readiness result\nShows gate-by-gate status", 0.8, 1.8, 2.35, 1.45, BLUE)
    add_panel(slide, "2. Scorecard", "Nine launch gates\nPass / fail / provisional\nThreshold-aware output", 3.35, 1.8, 2.2, 1.45, TEAL)
    add_panel(slide, "3. Unit Economics", "Projected revenue\nProjected cost\nOperating margin\nMargin buffer", 5.75, 1.8, 2.2, 1.45, GREEN)
    add_panel(slide, "4. Risk Logic", "Failed gates\nBlocking assumptions\nMitigation levers", 8.15, 1.8, 2.15, 1.45, AMBER)
    add_panel(slide, "5. Explanation Deck", "Reasoning\nFormula basis\nPlots and caveats\nExecutive summary", 10.5, 1.8, 2.0, 1.45, NAVY)
    add_callout(
        slide,
        "Delivery boundary",
        "Tomorrow's delivery is an MVP decision artifact. It is explicitly not a production deployment, not a live-market guarantee beyond the available inputs, and not an ML forecasting platform.",
        0.8,
        3.7,
        11.7,
        1.45,
        RED,
    )
    add_bullets(
        slide,
        [
            "The scope is anchored to the charter deliverables: readiness dashboard, intake flow, projected margin, risks, and mitigation recommendations.",
            "The plan deliberately cuts interface breadth before it cuts auditability or decision logic.",
        ],
        0.95,
        5.45,
        11.1,
        1.0,
        font_size=16,
    )
    add_footer(slide, "Scope for April 16, 2026", "2")

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    add_title(slide, "Why We Can Deliver This Quickly", "The project already has the minimum structure needed for a confidence-aware MVP")
    add_callout(slide, "Frozen business problem", "The core decision is fixed: does a market satisfy nine launch gates and support a 25% operating margin?", 0.8, 1.7, 5.6, 1.1, BLUE)
    add_callout(slide, "Frozen analytical base", "Phase 1 already locked the canonical historical base, dictionaries, lineage rules, join limits, and minimum viable subset.", 6.6, 1.7, 5.8, 1.1, GREEN)
    add_callout(slide, "Known blockers", "Kent-Leg, billed utilization, no-show denominator, pool logic, and cost granularity are already enumerated and do not need rediscovery.", 0.8, 3.0, 5.6, 1.1, AMBER)
    add_callout(slide, "Known exit strategy", "The plan already allows Tier 1 audited, Tier 2 assumption-backed, and Tier 3 manual-override outputs.", 6.6, 3.0, 5.8, 1.1, TEAL)
    add_callout(slide, "Usable demo input", "The intake example already contains enough structure to run a demonstrator flow across service lines, modes, counts, no-shows, and revenue assumptions.", 0.8, 4.3, 11.6, 1.1, NAVY)
    add_footer(slide, "Why the promise is credible", "3")

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    add_title(slide, "Dashboard MVP Experience", "Five panels cover the full decision flow without forcing stakeholders into details too early")
    add_panel(slide, "Intake Summary", "Market and organization\nService lines and modes\nWeekly trip counts\nNo-show inputs\nPricing assumptions", 0.8, 1.7, 3.7, 2.0, BLUE)
    add_panel(slide, "Readiness Scorecard", "Nine launch gates\nPass, fail, or provisional\nThreshold-aware output\nDistance from target", 4.8, 1.7, 3.1, 2.0, TEAL)
    add_panel(slide, "Margin & Economics", "Revenue\nCost\nOperating margin\nMargin buffer\nRevenue per Kent-Leg\nCost per road hour", 8.2, 1.7, 4.2, 2.0, GREEN)
    add_panel(slide, "Assumptions & Confidence", "Audited metrics\nAssumption-backed metrics\nFallback logic in use\nFormula branches where unresolved definitions still matter", 0.8, 4.0, 5.7, 1.7, AMBER)
    add_panel(slide, "Risk & Mitigation", "Failed gates\nBlocking assumptions\nLikely levers: pricing, mix, schedule density, no-show protection, and cost control", 6.8, 4.0, 5.6, 1.7, NAVY)
    add_callout(
        slide,
        "Design principle",
        "The dashboard is not just charts. It is a decision system: input, calculate, score, explain, and recommend.",
        0.8,
        5.95,
        11.6,
        0.7,
        TEAL,
    )
    add_footer(slide, "Target stakeholder experience", "4")

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    add_title(slide, "Analytical Architecture", "The MVP stays faithful to the documented five-module flow, but each module inherits confidence logic")
    add_panel(slide, "1. Demand", "Trips, miles, weekday shape, and mode mix become projected volume and Kent-Leg demand.", 0.8, 1.85, 3.7, 1.5, BLUE)
    add_panel(slide, "2. Contract Model", "Billing rules, no-show protection, payer structure, and quality vs filler become revenue protection and concentration logic.", 4.8, 1.85, 3.7, 1.5, TEAL)
    add_panel(slide, "3. Capacity & Scheduling", "Projected Kent-Legs, road-hour assumptions, and vehicle assumptions become utilization, road hours, and pool logic.", 8.8, 1.85, 3.6, 1.5, GREEN)
    add_panel(slide, "4. Cost", "Operating cost basis and cost-per-road-hour assumptions become projected cost.", 1.9, 3.85, 4.0, 1.45, AMBER)
    add_panel(slide, "5. Revenue & Margin", "Projected revenue, projected cost, and projected Kent-Legs become margin, buffer, and the final readiness state.", 7.0, 3.85, 4.4, 1.45, NAVY)

    add_callout(
        slide,
        "Implementation stance",
        "The dashboard will stay faithful to the documented modules, but every module will inherit the confidence logic already frozen in the execution plan.",
        0.8,
        5.7,
        11.6,
        0.8,
        BLUE,
    )
    add_footer(slide, "Five-module MVP backbone", "5")

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    add_title(slide, "Accuracy Safeguards", "The MVP remains credible by being explicit about what is known and what is still assumption-backed")
    add_callout(slide, "Tier 1 Audited", "Observed data and historically reconcilable logic. Safe for headline dashboard output.", 0.8, 1.8, 3.55, 1.2, GREEN)
    add_callout(slide, "Tier 2 Assumption-Backed", "Partially validated logic. Allowed in ranges, scenarios, and qualified explanations.", 4.6, 1.8, 3.55, 1.2, AMBER)
    add_callout(slide, "Tier 3 Manual Override", "Unresolved logic exposed explicitly. Cannot drive a fake hard claim.", 8.4, 1.8, 3.55, 1.2, RED)
    add_bullets(
        slide,
        [
            "Known ambiguities: Kent-Leg formula, billed-utilization formula, no-show denominator, pool logic, regional cost detail.",
            "Dashboard rule: unresolved items are surfaced as assumptions, branches, or provisional gates rather than normalized away.",
            "Decision outputs can therefore be Ready, Not Ready, or Insufficient Data without breaking the user experience.",
        ],
        0.9,
        3.45,
        11.1,
        1.55,
        font_size=16,
    )
    add_callout(
        slide,
        "Bottom line",
        "Accuracy comes from honest labeling, source precedence, and explicit fallback paths, not from pretending every formula is already final.",
        0.8,
        5.6,
        11.6,
        0.95,
        TEAL,
    )
    add_footer(slide, "Confidence-aware delivery", "6")

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    add_title(slide, "What Tomorrow's Delivery Includes", "A sharp bundle sized for decision support, not overbuild")
    add_bullets(
        slide,
        [
            "Dashboard MVP artifact using the existing intake and rules framework",
            "Nine-gate scorecard with pass / fail / provisional status",
            "Projected margin and unit-economics view",
            "Assumptions section with confidence labels",
            "Risk and mitigation explanation",
            "Short explanation deck with reasoning, formulas, and backing",
        ],
        0.9,
        1.85,
        5.4,
        3.2,
        font_size=17,
    )
    add_callout(
        slide,
        "Excluded from tomorrow's scope",
        "No new stakeholder interviews\nNo new external data collection\nNo production deployment stack\nNo full forecasting workbench",
        6.7,
        1.95,
        5.55,
        1.8,
        RED,
    )
    add_callout(
        slide,
        "Why this scope cut is correct",
        "The plan explicitly says to cut polish and extras before cutting auditability or decision logic. The real value is the recommendation engine, not speculative platform breadth.",
        6.7,
        4.15,
        5.55,
        1.45,
        NAVY,
    )
    add_footer(slide, "Scoped for a credible next-day MVP", "7")

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_background(slide, NAVY)
    add_title(slide, "Final Commitment", "What leadership can expect by Thursday, April 16, 2026", dark=True)
    add_callout(
        slide,
        "Committed outcome",
        "A presentable dashboard MVP and explanation deck that translates the existing data and rules into a usable market-readiness decision flow.",
        0.9,
        1.85,
        11.4,
        1.2,
        TEAL,
    )
    add_bullets(
        slide,
        [
            "The delivery will be professional, concise, and usable in an executive conversation under 10 minutes.",
            "It will be evidence-backed and confidence-aware.",
            "It will not hide unresolved assumptions.",
            "It will be aligned to the charter, the frozen MVP plan, and the canonical analytical base already built.",
        ],
        1.05,
        3.35,
        10.9,
        1.8,
        font_size=19,
        color=WHITE,
    )
    add_banner(slide, "PROMISE: DASHBOARD MVP + EXPLANATION DECK BY APRIL 16, 2026", 1.2, 5.55, 10.9, GREEN)
    add_footer(slide, "Ride YourWay Dashboard MVP Proposal", "8", dark=True)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build())
